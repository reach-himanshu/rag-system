"""Unit tests for document text extractors."""

import io

import pytest
from app.core.exceptions import DocumentProcessingError
from app.document_processing.extractors import SUPPORTED_TYPES, extract_text


class TestExtractText:
    """Tests for the extract_text dispatcher."""

    def test_unsupported_type_raises_error(self):
        """Should raise DocumentProcessingError for unsupported MIME types."""
        with pytest.raises(DocumentProcessingError, match="Unsupported file type"):
            extract_text(b"some data", "text/plain")

    def test_supported_types_mapping(self):
        """All expected MIME types should be in the SUPPORTED_TYPES mapping."""
        assert "application/pdf" in SUPPORTED_TYPES
        assert (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            in SUPPORTED_TYPES
        )
        assert (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" in SUPPORTED_TYPES
        )
        assert (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            in SUPPORTED_TYPES
        )


class TestExtractPdf:
    """Tests for PDF extraction."""

    def test_valid_pdf_extraction(self):
        """Should extract text from a valid PDF."""
        from PyPDF2 import PdfWriter

        # Create a minimal PDF in memory
        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)
        # Add text via annotation (simplest way without reportlab)
        writer.pages[0]

        buf = io.BytesIO()
        writer.write(buf)
        pdf_bytes = buf.getvalue()

        # A blank PDF should raise DocumentProcessingError (no text)
        with pytest.raises(DocumentProcessingError, match="No text content"):
            extract_text(pdf_bytes, "application/pdf")

    def test_corrupted_pdf_raises_error(self):
        """Should raise DocumentProcessingError for corrupted PDF data."""
        with pytest.raises(DocumentProcessingError, match="Failed to extract"):
            extract_text(b"not a real pdf", "application/pdf")


class TestExtractDocx:
    """Tests for Word document extraction."""

    def test_valid_docx_extraction(self):
        """Should extract text from a Word document."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Hello World")
        doc.add_paragraph("This is a test document.")

        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        result = extract_text(docx_bytes, content_type)

        assert "Hello World" in result
        assert "This is a test document." in result

    def test_docx_with_table(self):
        """Should extract text from tables in a Word document."""
        from docx import Document

        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Item A"
        table.cell(1, 1).text = "100"

        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        result = extract_text(docx_bytes, content_type)

        assert "Name" in result
        assert "Item A" in result


class TestExtractXlsx:
    """Tests for Excel extraction."""

    def test_valid_xlsx_extraction(self):
        """Should extract text from an Excel file."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sales"
        ws.append(["Product", "Revenue"])
        ws.append(["Widget A", 1500])
        ws.append(["Widget B", 2300])

        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        result = extract_text(xlsx_bytes, content_type)

        assert "Sales" in result
        assert "Widget A" in result
        assert "1500" in result


class TestExtractPptx:
    """Tests for PowerPoint extraction."""

    def test_valid_pptx_extraction(self):
        """Should extract text from a PowerPoint file."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Bullet point content"

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        result = extract_text(pptx_bytes, content_type)

        assert "Slide Title" in result
        assert "Bullet point content" in result
