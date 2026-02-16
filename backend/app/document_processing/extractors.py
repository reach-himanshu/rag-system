"""Text extraction from various document formats."""

import io
import logging

from app.core.exceptions import DocumentProcessingError

logger = logging.getLogger(__name__)

# Mapping of MIME types to extractor functions
SUPPORTED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}


def extract_text(file_bytes: bytes, content_type: str) -> str:
    """Extract text from a document based on its MIME type.

    Args:
        file_bytes: Raw file bytes.
        content_type: MIME type of the file.

    Returns:
        Extracted text content.

    Raises:
        DocumentProcessingError: If extraction fails or type is unsupported.
    """
    file_type = SUPPORTED_TYPES.get(content_type)
    if file_type is None:
        raise DocumentProcessingError(
            f"Unsupported file type: {content_type}. Supported: {', '.join(SUPPORTED_TYPES.keys())}"
        )

    extractors = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
    }

    try:
        text = extractors[file_type](file_bytes)
        if not text.strip():
            raise DocumentProcessingError("No text content found in document")
        logger.info("Extracted %d characters from %s file", len(text), file_type)
        return text
    except DocumentProcessingError:
        raise
    except Exception as e:
        raise DocumentProcessingError(f"Failed to extract text from {file_type}: {e}") from e


def _extract_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            pages.append(f"--- Page {i + 1} ---\n{text}")
    return "\n\n".join(pages)


def _extract_docx(file_bytes: bytes) -> str:
    """Extract text from a Word document, including tables."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip("| "):
                parts.append(row_text)

    return "\n".join(parts)


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract text from an Excel spreadsheet."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                parts.append(row_text)

    wb.close()
    return "\n".join(parts)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))

    return "\n\n".join(parts)
